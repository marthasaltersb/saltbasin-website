from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

files = [
    r'C:\Users\mbets\Downloads\go-to-market-strategy-mckinsey.pptx',
    r'C:\Users\mbets\Downloads\mckinsey-consulting-report.pptx',
    r'C:\Users\mbets\Downloads\data-strategy-roadmap.pptx',
    r'C:\Users\mbets\Downloads\project-status-report-heat-map.pptx',
    r'C:\Users\mbets\Downloads\data-infographic.pptx',
    r'C:\Users\mbets\Downloads\professional-monthly-scoreboard.pptx',
    r'C:\Users\mbets\Downloads\startup-executive-summary.pptx',
    r'C:\Users\mbets\Downloads\my-personal-dashboard.pptx',
    r'C:\Users\mbets\Downloads\floral-decision-matrix.pptx',
    r'C:\Users\mbets\Downloads\work-decision-tree.pptx',
    r'C:\Users\mbets\Downloads\employee-scorecard.pptx',
    r'C:\Users\mbets\Downloads\minimal-annual-company-review.pptx',
    r'C:\Users\mbets\Downloads\infographics.pptx',
    r'C:\Users\mbets\Downloads\phone-app.pptx',
]

for f in files:
    name = os.path.basename(f)
    print(f'\n{"="*60}')
    print(f'FILE: {name}')
    print('='*60)
    try:
        prs = Presentation(f)
        w = prs.slide_width.inches
        h = prs.slide_height.inches
        print(f'Dimensions: {w:.1f}" x {h:.1f}"  |  {len(prs.slides)} slides')

        for i, slide in enumerate(prs.slides[:2]):
            print(f'\n  -- Slide {i+1} ({len(slide.shapes)} shapes) --')
            # Get background color if any
            bg = slide.background
            fill = bg.fill
            try:
                if fill.type is not None:
                    print(f'  BG fill type: {fill.type}')
            except:
                pass

            for shape in slide.shapes:
                info = f'    shape_type={shape.shape_type} name="{shape.name}"'
                info += f' pos=({shape.left/914400:.2f}",{shape.top/914400:.2f}") size=({shape.width/914400:.2f}"x{shape.height/914400:.2f}")'

                # Try to get fill color
                try:
                    if hasattr(shape, 'fill'):
                        sf = shape.fill
                        if sf.type == 1:  # solid
                            rgb = sf.fore_color.rgb
                            info += f' fill=#{rgb}'
                except:
                    pass

                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        info += f'\n      TEXT: "{txt[:100]}"'
                        # Get font info from first paragraph
                        try:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                if run.font.size:
                                    info += f' [size={run.font.size.pt:.0f}pt'
                                if run.font.bold:
                                    info += ' bold'
                                try:
                                    if run.font.color.rgb:
                                        info += f' color=#{run.font.color.rgb}'
                                except:
                                    pass
                                info += ']'
                        except:
                            pass
                print(info)
    except Exception as e:
        print(f'  ERROR: {e}')

print('\nDone.')
